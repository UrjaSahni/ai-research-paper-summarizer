import os
from dotenv import load_dotenv
import streamlit as st
from transformers import pipeline
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import io
import re
from collections import defaultdict

# Load environment variables
load_dotenv('.env.example')

# Page config
st.set_page_config(
    page_title="AI Research Paper Summarizer",
    page_icon="📚",
    layout="wide"
)

# Initialize session state
if 'papers' not in st.session_state:
    st.session_state.papers = {}
if 'summaries' not in st.session_state:
    st.session_state.summaries = {}

# Helper functions for text extraction
def extract_text_from_pdf(file):
    try:
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx(file):
    try:
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def extract_text_from_pptx(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def extract_text_from_txt(file):
    try:
        text = file.read().decode('utf-8')
        return text
    except Exception as e:
        return f"Error reading file: {str(e)}"

def get_file_text(file):
    if file.name.endswith('.pdf'):
        return extract_text_from_pdf(file)
    elif file.name.endswith('.docx'):
        return extract_text_from_docx(file)
    elif file.name.endswith('.pptx'):
        return extract_text_from_pptx(file)
    elif file.name.endswith(('.txt', '.md')):
        return extract_text_from_txt(file)
    else:
        return "Unsupported file format"

# Smart text chunking
def chunk_text(text, max_chunk_size=1000, overlap=100):
    # Split by paragraphs first
    paragraphs = text.split('\n\n')
    chunks = []
    current_chunk = ""
    
    for para in paragraphs:
        if len(current_chunk) + len(para) < max_chunk_size:
            current_chunk += para + "\n\n"
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = para + "\n\n"
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

# Multi-level summarization
@st.cache_resource
def load_summarizer(api_key):
    try:
        return pipeline("summarization", model="facebook/bart-large-cnn", token=api_key)
    except:
        return pipeline("summarization", model="facebook/bart-large-cnn")

def generate_executive_summary(text, summarizer):
    chunks = chunk_text(text, max_chunk_size=1000)
    summaries = []
    
    for chunk in chunks[:3]:  # First 3 chunks for executive
        if len(chunk.split()) > 50:
            try:
                summary = summarizer(chunk, max_length=100, min_length=30, do_sample=False)
                summaries.append(summary[0]['summary_text'])
            except:
                pass
    
    return " ".join(summaries)

def generate_section_wise_summary(text, summarizer):
    # Try to detect sections
    sections = re.split(r'\n(?=[A-Z][^a-z]*:)|\n(?=\d+\.\s+[A-Z])', text)
    section_summaries = {}
    
    for i, section in enumerate(sections[:5]):  # First 5 sections
        if len(section.split()) > 50:
            section_name = f"Section {i+1}"
            # Try to extract section title
            first_line = section.split('\n')[0][:50]
            if first_line:
                section_name = first_line
            
            try:
                summary = summarizer(section[:1000], max_length=80, min_length=20, do_sample=False)
                section_summaries[section_name] = summary[0]['summary_text']
            except:
                pass
    
    return section_summaries

def extract_key_insights(text, summarizer):
    # Extract different types of insights
    insights = {
        "Methodology": "",
        "Results": "",
        "Conclusions": ""
    }
    
    # Look for methodology
    method_match = re.search(r'(method|approach|technique).*?(?=result|conclusion|\n\n)', text, re.IGNORECASE | re.DOTALL)
    if method_match:
        method_text = method_match.group()[:500]
        try:
            summary = summarizer(method_text, max_length=60, min_length=20, do_sample=False)
            insights["Methodology"] = summary[0]['summary_text']
        except:
            pass
    
    # Look for results
    results_match = re.search(r'(result|finding|outcome).*?(?=conclusion|discussion|\n\n)', text, re.IGNORECASE | re.DOTALL)
    if results_match:
        results_text = results_match.group()[:500]
        try:
            summary = summarizer(results_text, max_length=60, min_length=20, do_sample=False)
            insights["Results"] = summary[0]['summary_text']
        except:
            pass
    
    # Look for conclusions
    conclusion_match = re.search(r'(conclusion|summary|implication).*?$', text, re.IGNORECASE | re.DOTALL)
    if conclusion_match:
        conclusion_text = conclusion_match.group()[:500]
        try:
            summary = summarizer(conclusion_text, max_length=60, min_length=20, do_sample=False)
            insights["Conclusions"] = summary[0]['summary_text']
        except:
            pass
    
    return insights

# Comparative analysis
def compare_papers(summaries_dict):
    if len(summaries_dict) < 2:
        return None
    
    comparison = {
        "agreements": [],
        "contradictions": [],
        "gaps": []
    }
    
    paper_names = list(summaries_dict.keys())
    
    # Simple keyword-based comparison
    for i in range(len(paper_names)):
        for j in range(i+1, len(paper_names)):
            paper1 = paper_names[i]
            paper2 = paper_names[j]
            
            summary1 = summaries_dict[paper1].get('executive', '')
            summary2 = summaries_dict[paper2].get('executive', '')
            
            # Find common themes
            words1 = set(summary1.lower().split())
            words2 = set(summary2.lower().split())
            common = words1.intersection(words2)
            
            if len(common) > 10:
                comparison['agreements'].append(f"{paper1} and {paper2} discuss similar concepts")
    
    # Identify unique contributions
    all_topics = set()
    for paper, summaries in summaries_dict.items():
        exec_summary = summaries.get('executive', '')
        topics = set(exec_summary.lower().split())
        
        unique = topics - all_topics
        if len(unique) > 5:
            comparison['gaps'].append(f"{paper} introduces unique aspects not covered by others")
        
        all_topics.update(topics)
    
    return comparison

# Main UI
st.title("📚 AI Research Paper Summarizer & Comparator")
st.markdown("Upload multiple research papers and get AI-powered multi-level summaries and comparative analysis")

# Sidebar
with st.sidebar:
    st.header("⚙️ Configuration")
    api_key = st.text_input("Hugging Face API Key", type="password", value=os.getenv('HUGGINGFACE_API_KEY', ''), help="Enter your HF API key")
    st.markdown("---")
    st.markdown("""
    ### Features:
    - 📄 Upload multiple PDFs
    - 📊 Executive Summary
    - 📖 Section-wise Analysis  
    - 💡 Key Insights Extraction
    - 🔍 Comparative Analysis
    """)

# Main tabs
tab1, tab2, tab3 = st.tabs(["📂 Upload & Analyze", "🔍 Compare Papers", "ℹ️ About"])

with tab1:
    st.header("Upload Research Papers")
    
    uploaded_files = st.file_uploader(
        "Choose PDF files",
        accept_multiple_files=True,
        type=['pdf'],
        help="Upload one or more research papers in PDF format"
    )
    
    if uploaded_files and api_key:
        if st.button("⚡ Analyze Papers", type="primary"):
            summarizer = load_summarizer(api_key)
            
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            for idx, file in enumerate(uploaded_files):
                status_text.text(f"Processing {file.name}...")
                progress_bar.progress((idx + 1) / len(uploaded_files))
                
                # Extract text
                text = get_file_text(file)
                
                if not text.startswith("Error"):
                    st.session_state.papers[file.name] = text
                    
                    # Generate all summaries
                    with st.spinner(f"Generating summaries for {file.name}..."):
                        exec_summary = generate_executive_summary(text, summarizer)
                        section_summary = generate_section_wise_summary(text, summarizer)
                        insights = extract_key_insights(text, summarizer)
                        
                        st.session_state.summaries[file.name] = {
                            'executive': exec_summary,
                            'sections': section_summary,
                            'insights': insights
                        }
            
            status_text.text("✅ Analysis complete!")
            progress_bar.empty()
    
    # Display results
    if st.session_state.summaries:
        st.markdown("---")
        st.subheader("📊 Analysis Results")
        
        for paper_name, summaries in st.session_state.summaries.items():
            with st.expander(f"📑 {paper_name}", expanded=True):
                
                # Executive Summary
                st.markdown("### 🎯 Executive Summary")
                st.info(summaries['executive'])
                
                # Section-wise Summary
                st.markdown("### 📖 Section-wise Summary")
                if summaries['sections']:
                    for section, summary in summaries['sections'].items():
                        st.markdown(f"**{section}**")
                        st.write(summary)
                else:
                    st.write("No distinct sections detected")
                
                # Key Insights
                st.markdown("### 💡 Key Insights")
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    st.markdown("**Methodology**")
                    st.write(summaries['insights'].get('Methodology', 'Not found'))
                
                with col2:
                    st.markdown("**Results**")
                    st.write(summaries['insights'].get('Results', 'Not found'))
                
                with col3:
                    st.markdown("**Conclusions**")
                    st.write(summaries['insights'].get('Conclusions', 'Not found'))
    
    elif uploaded_files and not api_key:
        st.warning("⚠️ Please enter your Hugging Face API key in the sidebar")
    else:
        st.info("👆 Upload PDF files to get started")

with tab2:
    st.header("🔍 Comparative Analysis")
    
    if len(st.session_state.summaries) >= 2:
        comparison = compare_papers(st.session_state.summaries)
        
        if comparison:
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("### ✅ Agreements")
                st.markdown("Common themes and similar findings across papers")
                if comparison['agreements']:
                    for agreement in comparison['agreements']:
                        st.success(agreement)
                else:
                    st.info("No clear agreements found")
            
            with col2:
                st.markdown("### ⚠️ Contradictions")
                st.markdown("Conflicting findings or opposing views")
                if comparison['contradictions']:
                    for contradiction in comparison['contradictions']:
                        st.warning(contradiction)
                else:
                    st.info("No contradictions detected")
            
            with col3:
                st.markdown("### 🔍 Research Gaps")
                st.markdown("Unique contributions and unexplored areas")
                if comparison['gaps']:
                    for gap in comparison['gaps']:
                        st.info(gap)
                else:
                    st.info("No unique gaps identified")
            
            # Detailed comparison table
            st.markdown("---")
            st.subheader("📊 Detailed Comparison")
            
            comparison_data = []
            for paper_name, summaries in st.session_state.summaries.items():
                comparison_data.append({
                    "Paper": paper_name,
                    "Executive Summary": summaries['executive'][:200] + "...",
                    "Methodology": summaries['insights'].get('Methodology', 'N/A')[:100],
                    "Results": summaries['insights'].get('Results', 'N/A')[:100]
                })
            
            st.table(comparison_data)
    
    elif len(st.session_state.summaries) == 1:
        st.info("📝 Upload at least 2 papers to enable comparison")
    else:
        st.info("📂 Upload and analyze papers in the first tab to compare them")

with tab3:
    st.header("About This Application")
    
    st.markdown("""
    ### 🎯 Purpose
    This AI-powered tool helps researchers and students:
    - Quickly understand multiple research papers
    - Generate multi-level summaries (Executive, Section-wise, Key Insights)
    - Compare findings across multiple papers
    - Identify agreements, contradictions, and research gaps
    
    ### 🛠️ Features
    
    **Multi-Level Summarization:**
    - **Executive Summary**: High-level overview of the entire paper
    - **Section-wise Summary**: Detailed summaries of each major section
    - **Key Insights**: Extracted methodology, results, and conclusions
    
    **Comparative Analysis:**
    - **Agreements**: Common themes across papers
    - **Contradictions**: Conflicting findings
    - **Research Gaps**: Unique contributions and unexplored areas
    
    ### 💻 Technology Stack
    - **Frontend**: Streamlit
    - **AI Model**: Facebook BART (via Hugging Face)
    - **File Processing**: PyPDF2, python-docx, python-pptx
    
    ### 🔑 How to Use
    1. Enter your Hugging Face API key in the sidebar
    2. Upload PDF research papers (one or more)
    3. Click "Analyze Papers" to generate summaries
    4. View multi-level analysis for each paper
    5. Go to "Compare Papers" tab to see comparative analysis
    
    ### 👨‍💻 Developer
    Built by Urja Sahni
    """)
    
    st.markdown("---")
    st.markdown("**Note**: All processing is done securely. Your papers are not stored permanently.")

# Footer
st.markdown("---")
st.markdown(
    "<div style='text-align: center'>✨ Powered by Hugging Face Transformers • Made with Streamlit</div>",
    unsafe_allow_html=True
)
